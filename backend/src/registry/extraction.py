"""PPTX Inspection and Brand Token Extraction Service.

Uses python-pptx to analyze uploaded presentation decks and automatically extract
brand identity tokens (RGB color palettes, typography font families, and slide aspect ratio).
"""

from __future__ import annotations

import io
from collections import Counter
from typing import Any

from ..core.logging import get_logger

logger = get_logger("orchestrator.registry.extraction")

DEFAULT_PRIMARY = "#2563EB"
DEFAULT_SECONDARY = "#1E40AF"
DEFAULT_ACCENT = "#F59E0B"
DEFAULT_TYPOGRAPHY = "Inter"
DEFAULT_ASPECT_RATIO = "16:9"

SUPPORTED_FONTS = ["Inter", "JetBrains Mono", "Roboto", "Outfit", "Arial", "Calibri", "Helvetica"]


def extract_tokens_from_pptx(pptx_bytes: bytes) -> dict[str, Any]:
    """Inspect a .pptx file in memory and extract brand tokens."""
    try:
        import pptx
    except ImportError:
        logger.warning("python-pptx not installed, returning default brand tokens.")
        return _default_tokens()

    try:
        prs = pptx.Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        logger.error("Failed to parse PPTX file for token extraction: %s", exc)
        return _default_tokens()

    # 1. Aspect Ratio Detection
    aspect_ratio = DEFAULT_ASPECT_RATIO
    try:
        if prs.slide_height and prs.slide_width:
            ratio = prs.slide_width / prs.slide_height
            aspect_ratio = "16:9" if ratio > 1.5 else "4:3"
    except Exception as exc:
        logger.debug("Error computing aspect ratio: %s", exc)

    # 2. Typography & Color Harvesting
    fonts_counter: Counter[str] = Counter()
    colors_counter: Counter[str] = Counter()

    def _harvest_shape(shape: Any) -> None:
        try:
            # Check shape fill color
            if hasattr(shape, "fill") and shape.fill and hasattr(shape.fill, "fore_color"):
                color_obj = shape.fill.fore_color
                if hasattr(color_obj, "rgb") and color_obj.rgb is not None:
                    colors_counter[f"#{str(color_obj.rgb)}".upper()] += 1
        except Exception:
            pass

        try:
            # Check text font and text color
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, "font") and run.font:
                            if run.font.name:
                                fonts_counter[run.font.name] += 1
                            if (
                                hasattr(run.font, "color")
                                and run.font.color
                                and hasattr(run.font.color, "rgb")
                                and run.font.color.rgb is not None
                            ):
                                colors_counter[f"#{str(run.font.color.rgb)}".upper()] += 1
        except Exception:
            pass

    # Scan slide masters and layout slides
    try:
        for master in prs.slide_masters:
            for shape in master.shapes:
                _harvest_shape(shape)
            for layout in master.slide_layouts:
                for shape in layout.shapes:
                    _harvest_shape(shape)
        for slide in prs.slides:
            for shape in slide.shapes:
                _harvest_shape(shape)
    except Exception as exc:
        logger.debug("Error harvesting shapes: %s", exc)

    # 3. Resolve Typography
    typography = DEFAULT_TYPOGRAPHY
    if fonts_counter:
        top_font = fonts_counter.most_common(1)[0][0]
        # Try to match against our supported web fonts
        matched = False
        for supported in SUPPORTED_FONTS:
            if supported.lower() in top_font.lower() or top_font.lower() in supported.lower():
                typography = supported
                matched = True
                break
        if not matched and top_font:
            typography = top_font

    # 4. Resolve Color Palette (filter out white/black backgrounds)
    ignore_colors = {"#FFFFFF", "#FEFEFE", "#000000", "#010101", "#111827", "#1F2937", "#F3F4F6"}
    distinct_colors = [
        c for c, _ in colors_counter.most_common(20) if c not in ignore_colors and len(c) == 7
    ]

    primary_color = distinct_colors[0] if len(distinct_colors) >= 1 else DEFAULT_PRIMARY
    secondary_color = distinct_colors[1] if len(distinct_colors) >= 2 else DEFAULT_SECONDARY
    accent_color = distinct_colors[2] if len(distinct_colors) >= 3 else DEFAULT_ACCENT

    return {
        "primary_color": primary_color,
        "secondary_color": secondary_color,
        "accent_color": accent_color,
        "typography": typography,
        "aspect_ratio": aspect_ratio,
        "detected_fonts": [f for f, _ in fonts_counter.most_common(5)],
        "detected_colors": distinct_colors[:5],
    }


def _default_tokens() -> dict[str, Any]:
    return {
        "primary_color": DEFAULT_PRIMARY,
        "secondary_color": DEFAULT_SECONDARY,
        "accent_color": DEFAULT_ACCENT,
        "typography": DEFAULT_TYPOGRAPHY,
        "aspect_ratio": DEFAULT_ASPECT_RATIO,
        "detected_fonts": [DEFAULT_TYPOGRAPHY],
        "detected_colors": [DEFAULT_PRIMARY, DEFAULT_SECONDARY, DEFAULT_ACCENT],
    }
