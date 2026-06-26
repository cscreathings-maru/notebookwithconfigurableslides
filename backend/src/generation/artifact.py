"""Artifact inspection — read facts back out of the deck Presenton actually produced.

The consistency gate must judge the *produced* PPTX, not just the plan that drove it
(otherwise template-applied and slide-count are trivially true). This module opens the
generated PPTX bytes and reports the real slide count, slide titles (in order), and all
text — the inputs the checker compares against the profile's contract.
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field

from pptx import Presentation

from ..core.logging import get_logger

logger = get_logger("orchestrator.generation.artifact")


@dataclass(frozen=True)
class DeckFacts:
    """What the produced deck actually contains."""

    slide_count: int
    titles: tuple[str, ...] = field(default=())
    text: str = ""


def inspect_pptx(data: bytes) -> DeckFacts | None:
    """Parse a PPTX into DeckFacts, or None if the bytes are not a readable deck.

    A None result is itself a consistency signal (corrupt/empty artifact) handled by
    the checker — inspection never raises into the worker.
    """
    if not data:
        return None
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # not a valid OOXML package
        logger.warning("pptx_inspect_failed", extra={"error": str(exc)})
        return None

    titles: list[str] = []
    texts: list[str] = []
    for slide in prs.slides:
        title_text = ""
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip()
        titles.append(title_text)
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)

    return DeckFacts(
        slide_count=len(prs.slides),
        titles=tuple(titles),
        text="\n".join(texts),
    )
