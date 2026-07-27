"""Unit tests for PPTX token extraction service."""

import io
import pytest
from src.registry.extraction import extract_tokens_from_pptx, _default_tokens


def test_default_tokens() -> None:
    tokens = _default_tokens()
    assert tokens["primary_color"] == "#2563EB"
    assert tokens["secondary_color"] == "#1E40AF"
    assert tokens["accent_color"] == "#F59E0B"
    assert tokens["typography"] == "Inter"
    assert tokens["aspect_ratio"] == "16:9"


def test_extract_from_invalid_bytes() -> None:
    # Passing invalid bytes should gracefully fall back to default tokens
    tokens = extract_tokens_from_pptx(b"invalid pptx bytes")
    assert tokens["primary_color"] == "#2563EB"
    assert tokens["typography"] == "Inter"
    assert tokens["aspect_ratio"] == "16:9"


def test_extract_from_valid_pptx() -> None:
    try:
        import pptx
    except ImportError:
        pytest.skip("python-pptx not installed")

    # Create a minimal in-memory presentation
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    # Add a title shape with font name
    title = slide.shapes.title
    if title and title.has_text_frame:
        title.text = "Test Title"
        for p in title.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = "JetBrains Mono"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    tokens = extract_tokens_from_pptx(buf.getvalue())
    assert "primary_color" in tokens
    assert "secondary_color" in tokens
    assert "accent_color" in tokens
    assert "typography" in tokens
    assert "aspect_ratio" in tokens
    assert "detected_fonts" in tokens
    assert "detected_colors" in tokens
