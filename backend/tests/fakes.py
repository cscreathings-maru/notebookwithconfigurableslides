"""Test doubles for the ingestion slice.

These duck-type the collaborators the ingest service depends on (Open Notebook
client + object store) so upload->ready and failure paths are exercised without a
live engine, MinIO, Redis, or Arq.
"""

from __future__ import annotations

from collections.abc import Collection
from typing import Any

from src.engines.presenton import TemplateRegistration
from src.models import RegistrationStatus


class FakeObjectStore:
    """In-memory object store with the same surface as the MinIO store."""

    def __init__(self) -> None:
        self.objects: dict[str, bytes] = {}

    def tenant_key(self, *, tenant_id: str, project_id: str, source_id: str, filename: str) -> str:
        return f"{tenant_id}/{project_id}/sources/{source_id}/{filename}"

    def put_bytes(self, *, key: str, data: bytes, content_type: str) -> None:
        self.objects[key] = data

    def get_bytes(self, *, key: str) -> bytes:
        return self.objects[key]

    def presigned_get(self, *, key: str) -> str:
        return f"https://objectstore.test/{key}"


class FakeOpenNotebook:
    """Configurable fake of the Open Notebook client surface used by ingestion."""

    def __init__(
        self,
        *,
        notebook_id: str = "nb_fake",
        source_id: str = "src_fake",
        status_sequence: list[str] | None = None,
        analysis_ref: str = "analysis_fake",
        add_source_error: Exception | None = None,
        corpus: dict[str, str] | None = None,
    ) -> None:
        self.notebook_id = notebook_id
        self.source_id = source_id
        # Each get_source_status call pops the next status; defaults to ready.
        self._statuses = list(status_sequence or ["ready"])
        self.analysis_ref = analysis_ref
        self.add_source_error = add_source_error
        # {engine_source_ref: text}. Models the real engine's ONE GLOBAL INDEX -- every
        # project's content is visible to every query, which is exactly the condition
        # the caller-side scoping in search() has to survive.
        self.corpus = dict(corpus or {})
        self.searched_scopes: list[set[str]] = []
        self.calls: list[str] = []

    async def create_notebook(self, *, name: str, namespace: str) -> str:
        self.calls.append("create_notebook")
        return self.notebook_id

    async def add_source(
        self, *, notebook_id: str, uri: str, provider_config: dict[str, Any]
    ) -> str:
        self.calls.append("add_source")
        if self.add_source_error is not None:
            raise self.add_source_error
        return self.source_id

    async def get_source_status(self, *, source_id: str) -> str:
        self.calls.append("get_source_status")
        if len(self._statuses) > 1:
            return self._statuses.pop(0)
        return self._statuses[0]

    async def run_transformation(
        self, *, source_id: str, provider_config: dict[str, Any]
    ) -> str:
        self.calls.append("run_transformation")
        return self.analysis_ref

    async def search(
        self, *, allowed_source_refs: Collection[str], query: str
    ) -> list[dict[str, Any]]:
        """Serve from the global corpus, honouring the caller's allow-set.

        Mirrors the real client: the engine itself cannot scope, so anything not in
        `allowed_source_refs` must not come back.
        """
        self.calls.append("search")
        allowed = {str(r).strip() for r in allowed_source_refs if str(r).strip()}
        self.searched_scopes.append(allowed)
        if not allowed:
            return []
        if self.corpus:
            return [
                {"text": text, "source_ref": ref}
                for ref, text in self.corpus.items()
                if ref in allowed
            ]
        # No corpus configured: one canned snippet, attributed to an allowed source.
        return [{"text": "Revenue grew 12% YoY.", "source_ref": sorted(allowed)[0]}]


class FakeLlm:
    """Deterministic-structure LLM fake. Talking-point wording varies per call to
    prove that structure (sections/order) is fixed independently of the model."""

    def __init__(self) -> None:
        self._call = 0
        self.calls: list[str] = []
        self.chat_models: list[str | None] = []

    async def talking_points(self, *, section_ids, context, profile, provider_config):
        from src.outline.builder import LlmResult

        self._call += 1
        self.calls.append("talking_points")
        points = {sid: [f"point {sid} run{self._call}"] for sid in section_ids}
        return LlmResult(points_by_section=points, tokens_in=120, tokens_out=80)

    async def chat(
        self,
        *,
        system: str,
        user: str,
        provider_config: dict[str, Any],
        history: list[dict[str, str]] | None = None,
        temperature: float = 0.3,
        max_tokens: int = 1200,
        model_override: str | None = None,
    ):
        """Grounded completion fake. Returns a JSON question array when the prompt
        asks for questions, otherwise a plain grounded answer/summary."""
        from src.engines.llm import ChatAnswer

        self.calls.append("chat")
        self.chat_models.append(model_override or provider_config.get("model"))
        if "JSON" in system or "question" in system.lower():
            text = '["What drove revenue growth?", "What are the key risks?", "What is the outlook?"]'
        else:
            text = "Grounded overview: revenue grew 12% YoY based on the sources."
        return ChatAnswer(text=text, tokens_in=100, tokens_out=50)


def _pptx_from_markdown(content: str, slides_markdown: Any, n_slides: int) -> bytes:
    """Build a REAL PPTX from the orchestrator's params so the artifact-level
    consistency checker has a genuine deck to inspect: a title slide plus one slide
    per `## ` heading (title) with its `- ` bullets as body text. The total is padded
    to `n_slides` so the fake honors the requested count the way Presenton does.

    `slides_markdown` is Presenton's string[] (one block per slide); a bare string
    is tolerated for older callers."""
    import io

    from pptx import Presentation

    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = content or "Presentation"

    blocks = slides_markdown if isinstance(slides_markdown, list) else [slides_markdown or ""]
    current_body: Any = None
    for raw in "\n".join(blocks).splitlines():
        line = raw.strip()
        if line.startswith("## "):
            slide = prs.slides.add_slide(body_layout)
            slide.shapes.title.text = line[3:].strip()
            current_body = slide.placeholders[1].text_frame
            current_body.text = ""
        elif line.startswith("- ") and current_body is not None:
            para = current_body.add_paragraph()
            para.text = line[2:].strip()

    # Honor the requested slide count (Presenton pads within the profile's range).
    while len(prs.slides) < max(n_slides, 1):
        filler = prs.slides.add_slide(body_layout)
        filler.shapes.title.text = "Appendix"

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class FakePresenton:
    """Fake of the Presenton client surface (registry + generation)."""

    def __init__(self, *, ref_prefix: str = "tref", register_error: str | None = None) -> None:
        self.ref_prefix = ref_prefix
        # When set, registration reports a fallback with this reason (T-1.6).
        self.register_error = register_error
        self.registered: list[dict[str, Any]] = []
        self.generate_calls: list[dict[str, Any]] = []
        self.files: dict[str, bytes] = {}

    async def register_template(
        self,
        *,
        name: str,
        pptx_bytes: bytes | None = None,
        pptx_filename: str | None = None,
    ) -> TemplateRegistration:
        """Mirrors the real client: the PPTX *is* the brand (T-1.3).

        Without a deck there is nothing for the engine to derive colours, fonts or
        layouts from, so registration reports `fallback` rather than success.
        """
        self.registered.append({"name": name, "pptx_filename": pptx_filename})
        if self.register_error is not None:
            return TemplateRegistration.fallen_back(self.register_error)
        if not pptx_bytes:
            return TemplateRegistration.without_source()
        return TemplateRegistration(
            ref=f"{self.ref_prefix}_{name}",
            status=RegistrationStatus.registered,
            error=None,
        )

    async def generate(self, *, params: dict[str, Any]) -> dict[str, Any]:
        self.generate_calls.append(params)
        pid = f"pres_{len(self.generate_calls)}"
        export_as = str(params.get("export_as", "pptx")).lower()
        path = f"/app_data/{pid}.{export_as}"
        # Presenton returns one file in the requested format per generate call.
        if export_as == "pdf":
            self.files[path] = b"%PDF-1.4 fake-pdf"
        else:
            self.files[path] = _pptx_from_markdown(
                params.get("content", ""),
                params.get("slides_markdown", ""),
                int(params.get("n_slides", 0)),
            )
        return {"presentation_id": pid, "path": path, "edit_path": f"/edit/{pid}"}

    async def download(self, *, path: str) -> bytes:
        return self.files.get(path, b"")
